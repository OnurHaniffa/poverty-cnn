"""Build the poverty-cnn progress deck in the 'light premium' style.

Pure python-pptx. Off-white background, teal accent spine, eyebrow section
labels, navy titles with a teal rule, white content cards (with a faux soft
shadow), restyled transparent figures, and a footer with slide numbers.

Run: python scripts/07_build_deck.py   ->  results/poverty_cnn_progress_2026-05-21.pptx
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

# ---- palette ----
BG = RGBColor(0xF5, 0xF7, 0xFA); NAVY = RGBColor(0x1F, 0x3A, 0x5F)
SLATE = RGBColor(0x2B, 0x34, 0x40); TEAL = RGBColor(0x17, 0x8A, 0x7A)
TEAL_SOFT = RGBColor(0xE3, 0xF2, 0xEE); AMBER = RGBColor(0xE0, 0x92, 0x2F)
CARD = RGBColor(0xFF, 0xFF, 0xFF); BORDER = RGBColor(0xE1, 0xE7, 0xED)
GREY = RGBColor(0x7A, 0x84, 0x8F); SHADOW = RGBColor(0xE7, 0xEB, 0xF0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF); MIST = RGBColor(0xC9, 0xD6, 0xE3)
FONT = "Avenir Next"
D = "results/figures/diagrams/"; EG = "results/figures/eda/"
SW, SH, N = 10.0, 7.5, 11

prs = Presentation(); prs.slide_width = Inches(SW); prs.slide_height = Inches(SH)
BLANK = prs.slide_layouts[6]


def newslide():
    return prs.slides.add_slide(BLANK)


def rect(s, x, y, w, h, fill, line=None, lw=1.0, rounded=False, radius=0.05):
    shape = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    sp = s.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if rounded:
        try: sp.adjustments[0] = radius
        except Exception: pass
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb = line; sp.line.width = Pt(lw)
    sp.shadow.inherit = False
    return sp


def text(s, x, y, w, h, content, size=16, color=SLATE, bold=False,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space=8, line=1.06):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    paras = content if isinstance(content, list) else [content]
    for i, para in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(space); p.space_before = Pt(0); p.line_spacing = line
        segs = para if isinstance(para, list) else [(para, color, bold, size)]
        for seg in segs:
            t = seg[0]; col = seg[1] if len(seg) > 1 else color
            bd = seg[2] if len(seg) > 2 else bold; sz = seg[3] if len(seg) > 3 else size
            r = p.add_run(); r.text = t; r.font.name = FONT; r.font.size = Pt(sz)
            r.font.color.rgb = col; r.font.bold = bd
    return tb


def card(s, x, y, w, h):
    rect(s, x + 0.045, y + 0.06, w, h, SHADOW, rounded=True, radius=0.04)
    return rect(s, x, y, w, h, CARD, line=BORDER, lw=1.0, rounded=True, radius=0.04)


def pic(s, path, x, y, w):
    iw, ih = Image.open(path).size
    s.shapes.add_picture(path, Inches(x), Inches(y), Inches(w), Inches(w * ih / iw))


def BUL(t, size=15):
    return [("▸   ", TEAL, True, size), (t, SLATE, False, size)]


def chrome(s, eyebrow, title, idx, title_size=26):
    rect(s, 0, 0, SW, SH, BG)
    rect(s, 0, 0, 0.16, SH, TEAL)
    text(s, 0.6, 0.42, 8.8, 0.35, eyebrow.upper(), size=12.5, color=TEAL, bold=True)
    text(s, 0.55, 0.74, 9.2, 0.9, title, size=title_size, color=NAVY, bold=True, line=1.0)
    rect(s, 0.6, 1.52, 2.1, 0.045, TEAL)
    rect(s, 0.6, 7.0, 8.85, 0.013, BORDER)
    text(s, 0.6, 7.06, 5, 0.3, "poverty-cnn  ·  Onur Haniffa", size=10, color=GREY)
    text(s, 7.0, 7.06, 2.45, 0.3, f"{idx:02d} / {N:02d}", size=10, color=GREY, align=PP_ALIGN.RIGHT)


def headed_card(s, x, y, w, h, header, bullets, accent=TEAL):
    card(s, x, y, w, h)
    rect(s, x + 0.28, y + 0.3, 0.34, 0.06, accent)
    text(s, x + 0.28, y + 0.42, w - 0.56, 0.4, header.upper(), size=13, color=accent, bold=True)
    text(s, x + 0.28, y + 0.92, w - 0.56, h - 1.1, [BUL(b) for b in bullets], space=11, line=1.08)


def caption(s, t, y=6.35):
    text(s, 0.7, y, 8.6, 0.5, t, size=12.5, color=GREY, align=PP_ALIGN.CENTER)


# ===== Slide 0 — title =====
s = newslide()
rect(s, 0, 0, SW, SH, BG); rect(s, 0, 0, 0.16, SH, TEAL)
text(s, 0.7, 1.45, 8.6, 0.4, "ML/DL INTERNSHIP · SPRING 2026 · ACIBADEM MAAÜ", size=13, color=TEAL, bold=True)
text(s, 0.65, 1.95, 8.8, 1.7, "Predicting Village Wealth from Space", size=42, color=NAVY, bold=True, line=1.0)
rect(s, 0.72, 3.62, 2.6, 0.05, TEAL)
text(s, 0.7, 3.85, 8.5, 1.0, "A modern replication & fairness audit of Yeh et al. (2020), Nature Communications", size=17, color=SLATE, line=1.2)
rect(s, 0, 6.35, SW, 1.15, NAVY)
text(s, 0.7, 6.35, 8.6, 1.15, [[("Onur Haniffa", WHITE, True, 15),
    ("       ·       Advisor: Dr. Seda Nilgün Dumlu       ·       22 May 2026", MIST, False, 14)]],
    anchor=MSO_ANCHOR.MIDDLE)

# ===== Slide 1 — problem =====
s = newslide(); chrome(s, "Motivation", "The Problem & the Promise", 1)
headed_card(s, 0.55, 1.85, 4.35, 4.25, "The problem", [
    "Reliable poverty data is scarce, expensive, and infrequent",
    "Thinnest exactly where the need is greatest",
    "Decisions get made on stale or missing numbers"], accent=AMBER)
headed_card(s, 5.1, 1.85, 4.35, 4.25, "The promise", [
    "Satellites image everywhere on Earth, cheaply and repeatedly",
    "Can a neural network read household wealth from imagery?",
    "If yes: a scalable complement to ground surveys"], accent=TEAL)

# ===== Slide 2 — background + contributions =====
s = newslide(); chrome(s, "Background", "Where This Comes From", 2)
headed_card(s, 0.55, 1.85, 4.35, 4.25, "The lineage", [
    "Jean 2016 (Science): nightlights + transfer learning",
    "Yeh 2020 (Nature Comms): mean r² = 0.70 across 23 countries",
    "Aiken 2023 (IJCAI): fairness gaps — but only 10 countries"], accent=TEAL)
headed_card(s, 5.1, 1.85, 4.35, 4.25, "Our four contributions", [
    "Replicate Yeh in modern PyTorch",
    "Fairness audit across ALL 23 countries",
    "Uncertainty-aware fairness  (novel)",
    "Temporal fairness drift  (novel)"], accent=AMBER)

# ===== Slide 3 — big picture =====
s = newslide(); chrome(s, "Approach", "Two Data Sources, One Model", 3)
pic(s, D + "01_big_picture.png", 0.55, 2.15, 8.9)
caption(s, "Each village's GPS point pulls a satellite tile (x); its survey gives a wealth value (y); the CNN learns x → y.")

# ===== Slide 4 — wealth pipeline =====
s = newslide(); chrome(s, "Method · Wealth Index", "From a Survey to a Wealth Number", 4)
pic(s, D + "02_wealth_pipeline.png", 0.3, 2.45, 9.4)
caption(s, "Assets, not income. Pool every household across 23 countries, run PCA, take PC1 as the wealth axis, average to the village.", y=5.6)
caption(s, "PC1 explains 28.7% of total variance.", y=6.05)

# ===== Slide 5 — is PC1 valid =====
s = newslide(); chrome(s, "Method · Validity", "Is PC1 a Valid Wealth Index?", 5)
headed_card(s, 0.55, 1.85, 4.25, 4.5, "Why 28% is plenty", [
    "PC1 = 28.7% — but the benchmark is the 6.7% random floor, not 50%",
    "That is 4.3× the floor: one strongly dominant factor",
    "Binary survey data caps variance; 20–40% is the norm",
    "Variance = compactness, not validity",
    "Validated by the country ranking + urban/rural gap"])
card(s, 5.0, 1.85, 4.45, 4.5)
pic(s, EG + "01_pca_scree.png", 5.25, 2.2, 3.95)

# ===== Slide 6 — why PC1 is wealth =====
s = newslide(); chrome(s, "Method · Interpretation", "Why PC1 Is Wealth, Not PC2/PC3", 6)
headed_card(s, 0.55, 1.85, 4.25, 4.5, "Level vs contrast", [
    "A wealth factor makes ALL assets rise together (a 'level')",
    "PC1 is all-positive → that IS wealth",
    "PC2/PC3 are orthogonal contrasts (mixed signs)",
    "PC2 = rural/traditional vs urban/modern; PC3 = big durables",
    "Wealth is the #1 driver → it lands on PC1 by construction"])
card(s, 5.0, 1.85, 4.45, 4.5)
pic(s, EG + "02_pca_loadings.png", 5.3, 2.15, 3.85)

# ===== Slide 7 — GPS to tile =====
s = newslide(); chrome(s, "Method · Imagery", "From a GPS Point to a Tile", 7)
pic(s, D + "03_ee_bridge.png", 0.3, 2.3, 9.4)
caption(s, "The GPS coordinate selects the patch; Earth Engine builds a cloud-free 3-year median; 8 bands stack into one 224×224 tile.", y=5.55)

# ===== Slide 8 — preprocessing =====
s = newslide(); chrome(s, "Preprocessing", "Decisions & Honest Divergences", 8)
card(s, 0.55, 1.85, 8.9, 4.5)
text(s, 0.95, 2.2, 8.1, 3.9, [
    BUL("Ordinal categoricals (floor/wall/roof, water, toilet) collapsed to a binary 'improved/finished' flag — documented divergence"),
    BUL("Quality filter: drop households missing >30% of features — only ~30 of 355,445 dropped (99.99% complete)"),
    BUL("Special DHS codes handled (e.g. 'rooms' 96/97/98/99 → missing, not a fake 10-room house)"),
    BUL("Nightlights: DMSP-OLS (pre-2012) vs VIIRS (2012+) — different scales, documented"),
    BUL("One survey round per country → 13,634 clusters (Yeh pooled multiple rounds → 19,669)"),
], space=14, line=1.1)

# ===== Slide 9 — EDA =====
s = newslide(); chrome(s, "EDA", "What the Data Looks Like", 9)
card(s, 0.55, 1.8, 4.25, 3.95); pic(s, EG + "10_cluster_map.png", 0.95, 2.05, 3.45)
card(s, 5.0, 1.8, 4.45, 3.95); pic(s, EG + "07_urban_rural.png", 5.25, 2.5, 3.95)
for i, (num, lab, acc) in enumerate([("13,634", "georeferenced clusters", TEAL),
                                     ("28.7%", "PC1 variance", AMBER),
                                     ("1.16σ", "urban/rural gap", TEAL)]):
    x = 0.55 + i * 3.0
    card(s, x, 5.95, 2.85, 0.95)
    text(s, x, 6.08, 2.85, 0.5, num, size=24, color=acc, bold=True, align=PP_ALIGN.CENTER)
    text(s, x, 6.55, 2.85, 0.35, lab, size=11, color=GREY, align=PP_ALIGN.CENTER)

# ===== Slide 10 — status =====
s = newslide(); chrome(s, "Status", "Status & Plan", 10)
headed_card(s, 0.55, 1.85, 2.85, 3.6, "Done", [
    "DHS → wealth index, all 23 countries",
    "355,445 households → 13,634 villages",
    "Full EDA"], accent=TEAL)
headed_card(s, 3.6, 1.85, 2.85, 3.6, "In progress", [
    "Satellite tiles via Earth Engine",
    "~50% extracted",
    "Finishing in ~1–2 days"], accent=AMBER)
headed_card(s, 6.6, 1.85, 2.85, 3.6, "Next", [
    "Pair tiles ↔ wealth",
    "Train 8-channel ResNet-18",
    "Replication r² → fairness → extensions"], accent=NAVY)
card(s, 0.55, 5.65, 8.9, 0.95)
text(s, 0.55, 5.78, 8.9, 0.7, [[("Target  ", NAVY, True, 14), ("mean-of-folds r² ≥ 0.60   ", SLATE, False, 14),
    ("·   Deadline  ", NAVY, True, 14), ("30 May 2026", SLATE, False, 14)]],
    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

OUT = "results/poverty_cnn_progress_2026-05-21.pptx"
prs.save(OUT)
print("saved", OUT, "with", len(prs.slides._sldIdLst), "slides")
