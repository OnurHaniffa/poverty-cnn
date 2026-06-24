"""Final lab-presentation deck (~17 slides) in the 'light premium' style.

Teaching-first: each concept Why -> What -> Show -> So-what. Reuses the design
system from 07_build_deck.py. Run: python scripts/30_build_final_deck.py
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

BG=RGBColor(0xF5,0xF7,0xFA); NAVY=RGBColor(0x1F,0x3A,0x5F); SLATE=RGBColor(0x2B,0x34,0x40)
TEAL=RGBColor(0x17,0x8A,0x7A); TEAL_SOFT=RGBColor(0xE3,0xF2,0xEE); AMBER=RGBColor(0xE0,0x92,0x2F)
AMBER_SOFT=RGBColor(0xFB,0xEC,0xD6); RED=RGBColor(0xC0,0x50,0x4D); RED_SOFT=RGBColor(0xF6,0xE4,0xE3)
CARD=RGBColor(0xFF,0xFF,0xFF); BORDER=RGBColor(0xE1,0xE7,0xED); GREY=RGBColor(0x7A,0x84,0x8F)
SHADOW=RGBColor(0xE7,0xEB,0xF0); WHITE=RGBColor(0xFF,0xFF,0xFF); MIST=RGBColor(0xC9,0xD6,0xE3)
FONT="Avenir Next"
D="results/figures/diagrams/"; EG="results/figures/eda/"; TV="results/figures/tile_viz/"
F="results/figures/fairness/"; TC="results/figures/teaching/"
SW,SH,N=10.0,7.5,19
prs=Presentation(); prs.slide_width=Inches(SW); prs.slide_height=Inches(SH); BLANK=prs.slide_layouts[6]

def ns(): return prs.slides.add_slide(BLANK)
def rect(s,x,y,w,h,fill,line=None,lw=1.0,rounded=False,radius=0.05):
    sh=MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    sp=s.shapes.add_shape(sh,Inches(x),Inches(y),Inches(w),Inches(h))
    if rounded:
        try: sp.adjustments[0]=radius
        except Exception: pass
    sp.fill.solid(); sp.fill.fore_color.rgb=fill
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb=line; sp.line.width=Pt(lw)
    sp.shadow.inherit=False; return sp
def text(s,x,y,w,h,content,size=16,color=SLATE,bold=False,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,space=8,line=1.06):
    tb=s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); tf=tb.text_frame
    tf.word_wrap=True; tf.vertical_anchor=anchor
    for i,para in enumerate(content if isinstance(content,list) else [content]):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment=align; p.space_after=Pt(space); p.space_before=Pt(0); p.line_spacing=line
        for seg in (para if isinstance(para,list) else [(para,color,bold,size)]):
            t=seg[0]; col=seg[1] if len(seg)>1 else color; bd=seg[2] if len(seg)>2 else bold; sz=seg[3] if len(seg)>3 else size
            r=p.add_run(); r.text=t; r.font.name=FONT; r.font.size=Pt(sz); r.font.color.rgb=col; r.font.bold=bd
    return tb
def card(s,x,y,w,h):
    rect(s,x+0.045,y+0.06,w,h,SHADOW,rounded=True,radius=0.04)
    return rect(s,x,y,w,h,CARD,line=BORDER,lw=1.0,rounded=True,radius=0.04)
def pic(s,path,x,y,w):
    iw,ih=Image.open(path).size; s.shapes.add_picture(path,Inches(x),Inches(y),Inches(w),Inches(w*ih/iw))
def picture_h(path,w):
    iw,ih=Image.open(path).size; return w*ih/iw
def BUL(t,size=15,c=SLATE): return [("▸   ",TEAL,True,size),(t,c,False,size)]
def chrome(s,eyebrow,title,idx,title_size=26):
    rect(s,0,0,SW,SH,BG); rect(s,0,0,0.16,SH,TEAL)
    text(s,0.6,0.42,8.8,0.35,eyebrow.upper(),size=12.5,color=TEAL,bold=True)
    text(s,0.55,0.74,9.2,0.9,title,size=title_size,color=NAVY,bold=True,line=1.0)
    rect(s,0.6,1.52,2.1,0.045,TEAL); rect(s,0.6,7.0,8.85,0.013,BORDER)
    text(s,0.6,7.06,5,0.3,"Predicting Village Wealth from Space  ·  Onur Haniffa",size=10,color=GREY)
    text(s,7.0,7.06,2.45,0.3,f"{idx:02d} / {N:02d}",size=10,color=GREY,align=PP_ALIGN.RIGHT)
def headed_card(s,x,y,w,h,header,bullets,accent=TEAL):
    card(s,x,y,w,h); rect(s,x+0.28,y+0.3,0.34,0.06,accent)
    text(s,x+0.28,y+0.42,w-0.56,0.4,header.upper(),size=13,color=accent,bold=True)
    text(s,x+0.28,y+0.92,w-0.56,h-1.1,[BUL(b) for b in bullets],space=11,line=1.08)
def caption(s,t,y=6.35,col=GREY,size=12.5):
    text(s,0.7,y,8.6,0.5,t,size=size,color=col,align=PP_ALIGN.CENTER)
def chip(s,x,y,w,num,lab,acc=TEAL,h=1.15,nsz=27):
    card(s,x,y,w,h)
    text(s,x,y+0.14,w,0.6,num,size=nsz,color=acc,bold=True,align=PP_ALIGN.CENTER)
    text(s,x,y+0.72,w,0.35,lab,size=11,color=GREY,align=PP_ALIGN.CENTER)

# ===== S0 title =====
s=ns(); rect(s,0,0,SW,SH,BG); rect(s,0,0,0.16,SH,TEAL)
text(s,0.7,1.35,8.6,0.4,"ML / DL INTERNSHIP  ·  ACIBADEM MAAÜ  ·  2026",size=13,color=TEAL,bold=True)
text(s,0.65,1.85,8.9,1.8,"Predicting Village Wealth\nfrom Space",size=40,color=NAVY,bold=True,line=1.02)
rect(s,0.72,3.78,2.6,0.05,TEAL)
text(s,0.7,4.0,8.6,1.0,"A modern replication — and fairness & uncertainty audit — of Yeh et al. (2020), Nature Communications",size=16.5,color=SLATE,line=1.25)
rect(s,0,6.35,SW,1.15,NAVY)
text(s,0.7,6.35,8.6,1.15,[[("Onur Haniffa",WHITE,True,15),("        ·        Advisor: Dr. Seda Nilgün Dumlu        ·        June 2026",MIST,False,14)]],anchor=MSO_ANCHOR.MIDDLE)

# ===== S1 motivation =====
s=ns(); chrome(s,"Motivation","Why Predict Poverty from Space?",1)
headed_card(s,0.55,1.85,4.35,4.3,"The problem",[
    "Reliable poverty data is scarce, expensive, and infrequent",
    "Household surveys cost millions and run every ~5 years",
    "Thinnest exactly where the need is greatest"],accent=AMBER)
headed_card(s,5.1,1.85,4.35,4.3,"The promise",[
    "Satellites image everywhere on Earth — free and repeatedly",
    "Can a neural network read village wealth from imagery?",
    "If yes: a scalable complement to ground surveys"],accent=TEAL)
caption(s,"The whole question: can a free satellite image stand in for an expensive survey?",y=6.4,size=13)

# ===== S2 the idea =====
s=ns(); chrome(s,"The Idea","One Model, Two Inputs",2)
pic(s,D+"01_big_picture.png",0.55,2.2,8.9)
caption(s,"A village's GPS pulls a satellite tile (x); its survey gives a wealth value (y); the network learns x → y.",y=6.25)
caption(s,"This talk answers two things: where does the wealth number come from, and does it work — for whom?",y=6.62)

# ===== S3 DHS =====
s=ns(); chrome(s,"Data · Ground Truth","The Answer Key: DHS Surveys",3)
headed_card(s,0.55,1.85,4.35,4.3,"What DHS gives us",[
    "Gold-standard household surveys (USAID-funded)",
    "Unit = a cluster ≈ a village (~25 households)",
    "A GPS coordinate for every cluster"],accent=TEAL)
headed_card(s,5.1,1.85,4.35,4.3,"The deliberate twist",[
    "DHS shifts each GPS point for privacy: up to 2 km (urban), 5 km (rural)",
    "So the true village is somewhere nearby — not exactly on the point",
    "That is WHY our tile is a wide 6.72 km box: big enough to still contain it"],accent=AMBER)
caption(s,"The survey is our ‘answer key’ — the model is trained and graded against it.",y=6.4,size=13)

# ===== S4 no wealth column =====
s=ns(); chrome(s,"Data · The Target","There Is No ‘Wealth’ Column",4)
headed_card(s,0.55,1.85,4.25,4.4,"The problem",[
    "DHS never asks ‘how rich are you?’",
    "It records asset checkboxes: electricity? TV? fridge? car?",
    "Plus finished floor / wall / roof, water, toilet, rooms",
    "We must turn ~15 yes/no answers into ONE wealth number"])
card(s,5.0,1.85,4.45,4.4); pic(s,EG+"03_asset_prevalence.png",5.35,2.25,3.75)
caption(s,"Asset ownership across our 355,000 households — the raw material for the wealth index.",y=6.45,size=12)

# ===== S5 PCA =====
s=ns(); chrome(s,"Method · Wealth Index","PCA: Finding the Hidden Wealth Axis",5)
headed_card(s,0.55,1.85,4.25,4.5,"The one ruler",[
    "PCA finds the single mix of assets that best explains who owns MORE",
    "That axis (PC1) is wealth — every asset loads positive (a ‘more-of-everything’ axis)",
    "Average it to the village = our target y",
    "‘Only 28%?’ — that is 4× the random floor; validity is in the loadings, not the %"])
card(s,5.0,1.85,4.45,4.5); pic(s,EG+"02_pca_loadings.png",5.35,2.15,3.75)
caption(s,"Like a credit-score formula — discovered from the data, not hand-set.",y=6.5,size=12)

# ===== S6 the 8 bands (money shot) =====
s=ns(); chrome(s,"Data · The Input","Eight Ways to See a Village",6)
# urban column
text(s,0.55,1.78,4.35,0.3,"URBAN — NAIROBI, KENYA",size=12.5,color=TEAL,bold=True,align=PP_ALIGN.CENTER)
pic(s,TV+"nairobi_all_bands.png",0.55,2.12,4.35)
pic(s,TV+"nairobi_composites.png",0.55,2.12+picture_h(TV+"nairobi_all_bands.png",4.35)+0.12,4.35)
# rural column
text(s,5.1,1.78,4.35,0.3,"RURAL — TURKANA, KENYA",size=12.5,color=AMBER,bold=True,align=PP_ALIGN.CENTER)
pic(s,TV+"turkana_all_bands.png",5.1,2.12,4.35)
pic(s,TV+"turkana_composites.png",5.1,2.12+picture_h(TV+"turkana_all_bands.png",4.35)+0.12,4.35)
caption(s,"8 channels = 8 ‘senses’: visible colour, infrared (vegetation & rooftops), thermal (heat), night-lights (electrification). The CNN reads all 8 together and outputs one wealth number.",y=6.42,size=11.5)

# ===== S7 the model =====
s=ns(); chrome(s,"Method · The Model","The Model, in 20 Seconds",7)
card(s,0.55,1.95,8.9,3.0)
text(s,1.0,2.35,8.0,2.3,[
    BUL("A ResNet-18 — a standard, well-understood convolutional network"),
    BUL("Two tweaks: the input is rebuilt for 8 channels (not 3), and the output is one number (not categories)"),
    BUL("Trained from scratch on the satellite tiles"),
],space=15,line=1.12)
card(s,0.55,5.15,8.9,1.2)
text(s,0.55,5.15,8.9,1.2,[[("The science is in the data and the honest evaluation — not the architecture.",NAVY,True,15)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)

# ===== S8 cross-country =====
s=ns(); chrome(s,"Method · The Honest Test","Testing on Countries It Has Never Seen",8)
headed_card(s,0.55,1.85,4.35,4.3,"Why not a random split?",[
    "Split villages randomly and the model can MEMORISE a country",
    "Same country in train AND test = it cheats — a flattering but fake score",
    "Useless for the real question: does it generalise to a NEW place?"],accent=AMBER)
headed_card(s,5.1,1.85,4.35,4.3,"Cross-country validation",[
    "Train on ~18 countries, test on 5 it has NEVER seen",
    "Rotate so every country is the test set exactly once",
    "Every number we report is on unseen countries"],accent=TEAL)
caption(s,"Like an exam with questions you never studied — the only fair test.",y=6.4,size=13)

# ===== S9 metrics =====
s=ns(); chrome(s,"Method · Measuring Success","How Do We Know If It’s Good?",9)
text(s,0.6,1.78,8.8,0.4,"Three metrics — because they answer different questions.",size=14,color=SLATE,bold=True)
pic(s,TC+"metrics_example.png",0.5,2.25,9.0)

# ===== S10 replication =====
s=ns(); chrome(s,"Results · Replication","Does It Work?  Yes — at Benchmark Level.",10)
chip(s,0.55,1.95,2.1,"0.78","Pearson r (ours)",TEAL)
chip(s,2.85,1.95,2.1,"0.78","= WILDS benchmark",GREY)
chip(s,5.15,1.95,2.1,"0.55","worst-group r",TEAL)
chip(s,7.45,1.95,2.0,"0.45","WILDS worst-grp",GREY)
card(s,0.55,3.45,8.9,2.9)
text(s,1.0,3.85,8.0,2.2,[
    BUL("In Pearson r — the metric WILDS PovertyMap reports — we MATCH the published benchmark (0.78)…"),
    BUL("…and BEAT its worst-group score (0.55 vs 0.45), on the same leave-country-out protocol"),
    BUL("Pooled r² 0.61 on unseen countries; the ‘gap to Yeh’ was an r²-framing of the from-scratch penalty"),
    BUL("Over 2× a naïve floor — it genuinely reads wealth from space, on countries it has never seen"),
],space=13,line=1.12)

# ===== S11 nightlights =====
s=ns(); chrome(s,"Results · What Drives It","What Carries the Signal? Night-lights.",11)
chip(s,1.0,2.0,2.4,"0.58","night-lights only",AMBER,nsz=30)
chip(s,3.8,2.0,2.4,"0.52","all 8 channels",TEAL,nsz=30)
chip(s,6.6,2.0,2.4,"0.46","daytime only",GREY,nsz=30)
card(s,0.55,3.55,8.9,2.75)
text(s,1.0,3.95,8.0,2.0,[
    BUL("Night-lights ALONE roughly match the full 8-channel model"),
    BUL("Lights = electrification = the single strongest village-wealth signal"),
    BUL("Daytime imagery adds little here (it is data-hungry, and we are data-limited)"),
    BUL("Foreshadows the limit: if the signal is mostly lights, what happens where there are none?"),
],space=13,line=1.12)

# ===== S12 audit urban/rural =====
s=ns(); chrome(s,"Results · The Audit (1/2)","It Serves Cities Better Than Villages",12)
card(s,0.55,1.85,8.9,3.55); pic(s,F+"01_urban_rural.png",0.95,2.15,8.1)
card(s,0.55,5.6,8.9,1.0)
text(s,0.55,5.6,8.9,1.0,[[("Remember Nairobi vs Turkana?  ",NAVY,True,14),("Equal absolute accuracy (MAE), but the model RANKS rural villages worse (Spearman 0.41 vs 0.54) — in 21 of 23 countries.",SLATE,False,14)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)

# ===== S13 audit poorest =====
s=ns(); chrome(s,"Results · The Audit (2/2)","Confidently Wrong About the Poorest",13)
card(s,0.55,1.85,4.3,4.0); pic(s,F+"05_calibration.png",1.35,2.15,2.7)
card(s,5.05,1.85,4.4,4.0); pic(s,F+"06_bias_by_wealth.png",5.25,2.6,4.0)
caption(s,"The model compresses toward the average: it predicts the POOREST as richer than they are, with the biggest errors exactly there. Targeting the poorest 20%, it finds only ~23% of them.",y=6.05,size=12)

# ===== S14 the punchline =====
s=ns(); chrome(s,"The Punchline","A Fundamental Limit",14)
boxes=[("CONFIDENTLY WRONG","The audit: largest, most biased errors fall on the poorest villages",RED,RED_SOFT),
       ("UNDETECTABLE","No uncertainty method — ensembles, MC-dropout, learned variance — flags it",AMBER,AMBER_SOFT),
       ("UNFIXABLE","Reweighting the loss toward the poor doesn’t repair it — it just destabilises training",TEAL,TEAL_SOFT)]
for i,(h,b,acc,soft) in enumerate(boxes):
    x=0.55+i*3.02
    rect(s,x,1.95,2.85,2.5,soft,rounded=True,radius=0.05)
    rect(s,x+0.25,2.2,0.32,0.06,acc)
    text(s,x+0.25,2.32,2.4,0.5,h,size=12.5,color=acc,bold=True)
    text(s,x+0.25,2.78,2.4,1.5,b,size=12.5,color=SLATE,line=1.12)
    if i<2: text(s,x+2.78,2.9,0.3,0.5,"→",size=22,color=GREY,bold=True)
card(s,0.55,4.75,8.9,1.7)
text(s,0.9,4.98,8.2,1.3,[[("Because the signal isn’t there.  ",NAVY,True,16),("Night-lights are uniformly dark below the poverty line — the satellite cannot separate the extreme poor.",SLATE,False,14.5)],
    [("Confirmed FIVE ways: ",NAVY,True,14),("confidently-wrong → undetectable → unfixable by reweighting (LDS & Balanced-MSE) → unfixable by 3× data → and it recurs out-of-distribution (next slides).",SLATE,False,14)]],space=8,line=1.14)

# ===== S15 full study confirms =====
s=ns(); chrome(s,"Results · Scale","The Full Study Confirms It (3× Data)",15)
chip(s,0.55,1.95,2.85,"0.61","pooled r² (full)",TEAL)
chip(s,3.6,1.95,2.85,"+0.59","poorest-bias",AMBER)
chip(s,6.6,1.95,2.85,"≈","fairness unchanged",GREY)
card(s,0.55,3.4,4.35,2.95)
text(s,0.85,3.72,3.8,2.5,[
    BUL("3× the data (multi-round): accuracy rises 0.52 → 0.61"),
    BUL("But the poorest-bias barely moves (+0.62 → +0.59)"),
    BUL("More data helps accuracy — NOT equity"),
],space=12,line=1.13)
card(s,5.05,3.4,4.4,2.95); pic(s,TC+"data_scaling_curve.png",5.3,3.7,3.9)

# ===== S16 OOD capstone =====
s=ns(); chrome(s,"Results · OOD Capstone","Where Does It Break? Unseen Countries.",16)
card(s,0.55,1.8,8.9,4.0); pic(s,TC+"ood_where_it_breaks.png",2.05,2.05,5.9)
card(s,0.55,5.95,8.9,0.85)
text(s,0.7,5.95,8.6,0.85,[[("Frozen model, 6 never-seen countries.  ",NAVY,True,13),("Ranking TRANSFERS (Gabon 0.81 beats home) — but it BREAKS on rich South Africa (r² −1.7): the same flatten-the-extremes flaw, now on the rich tail. And more data does NOT help OOD.",SLATE,False,13.5)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)

# ===== S17 evaluation honesty =====
s=ns(); chrome(s,"Method · Honesty","Is Our Evaluation Honest? We Checked.",17)
chip(s,0.55,1.95,2.85,"+20%","what naïve CV inflates",AMBER)
chip(s,3.6,1.95,2.85,"23/23","survive FDR control",TEAL)
chip(s,6.6,1.95,2.85,"0.61","[0.54–0.66] honest CI",TEAL)
card(s,0.55,3.4,8.9,2.95)
text(s,1.0,3.74,8.0,2.5,[
    BUL("Random-split CV scores 0.73 vs our country-blocked 0.61 — a +20% leakage inflation we REFUSE to claim"),
    BUL("No memorisation: train r² 0.70–0.80 (not ~1.0); learning curve early-stopped at best validation"),
    BUL("Whole-country bootstrap CIs + Benjamini-Hochberg FDR across all 23 per-country tests"),
    BUL("Residuals carry local spatial structure — unmodelled signal, NOT leakage (folds are country-disjoint)"),
],space=12,line=1.13)

# ===== S18 contributions =====
s=ns(); chrome(s,"Contributions","What’s New Here",18)
headed_card(s,0.55,1.85,4.35,4.3,"Solid foundation",[
    "Faithful modern replication — matches the WILDS benchmark in Pearson r",
    "The most comprehensive fairness + uncertainty audit on this 23-country data",
    "Adds temporal + out-of-distribution external validity (frozen models, 6 new countries)"],accent=TEAL)
headed_card(s,5.1,1.85,4.35,4.3,"Most novel",[
    "Equity-framed uncertainty: even calibrated UQ cannot flag the failure on the poorest",
    "A fundamental-limit thesis confirmed FIVE independent ways — apparently unstated in the literature"],accent=AMBER)
caption(s,"Honest positioning: this is a rigorous AUDIT, not a new architecture — and we don’t claim a new state of the art.",y=6.4,size=13)

# ===== S19 takeaway =====
s=ns(); rect(s,0,0,SW,SH,NAVY); rect(s,0,0,0.16,SH,TEAL)
text(s,0.7,1.5,8.6,0.4,"THE ONE THING TO REMEMBER",size=13,color=TEAL,bold=True)
text(s,0.7,2.2,8.6,3.0,[[("Satellite poverty maps are great for ",WHITE,False,23),("ranking regions",TEAL,True,23),(" — not for last-mile targeting of the extreme poor — and standard uncertainty tools can’t warn you when they’re wrong about the neediest.",WHITE,False,23)]],line=1.25)
rect(s,0.72,5.55,2.6,0.05,TEAL)
text(s,0.7,5.8,8.6,0.5,"Thank you  ·  Questions?",size=18,color=MIST,bold=True)

OUT="results/poverty_cnn_final_deck.pptx"
prs.save(OUT); print("saved",OUT,"with",len(prs.slides._sldIdLst),"slides")
