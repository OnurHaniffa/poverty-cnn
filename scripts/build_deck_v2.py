"""Poverty-CNN deck v2 — rebuilt in the 'engineering blueprint' style (graph-paper bg, faint
gears, cyan blueprint node-lines, cyan chevrons, teal-blue bold headlines, green number-circles,
dark-teal emphasis boxes). All motifs are drawn here; imagery is our own satellite tiles/figures.

Run on the Mac:  /opt/homebrew/bin/python3.13 scripts/build_deck_v2.py
"""
from pathlib import Path
import math
from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---- palette (sampled from the EMG deck) ----
BG      = RGBColor(0xED,0xED,0xEB)
TEAL    = RGBColor(0x22,0x6E,0x8C)   # headlines
TEAL_D  = RGBColor(0x12,0x4A,0x63)   # darker
CYAN    = RGBColor(0x16,0xCE,0xE2)   # chevrons / node dots
GREEN   = RGBColor(0x1B,0xAE,0x8B)   # number circles
DARKBOX = RGBColor(0x10,0x52,0x74)   # emphasis box
INK     = RGBColor(0x2B,0x30,0x36)   # body
GREY    = RGBColor(0x8A,0x93,0x9A)
WHITE   = RGBColor(0xFF,0xFF,0xFF)
LINE    = RGBColor(0x23,0x29,0x2E)
HEAD="Avenir Next"; BODY="Avenir Next"
SW,SH=13.333,7.5
ASSET=Path("results/figures/assets"); ASSET.mkdir(parents=True,exist_ok=True)
TV="results/figures/tile_viz/"


# ---------- background PNG (grid + faint gears) ----------
def make_bg():
    W,H=1920,1080
    im=Image.new("RGB",(W,H),(0xED,0xED,0xEB)); d=ImageDraw.Draw(im,"RGBA")
    step=34
    for x in range(0,W,step): d.line([(x,0),(x,H)],fill=(0,0,0,16),width=1)
    for y in range(0,H,step): d.line([(0,y),(W,y)],fill=(0,0,0,16),width=1)
    def gear(cx,cy,r,teeth=14,col=(0,0,0,16)):
        for i in range(teeth):
            a=2*math.pi*i/teeth
            d.line([(cx+r*0.7*math.cos(a),cy+r*0.7*math.sin(a)),
                    (cx+r*math.cos(a),cy+r*math.sin(a))],fill=col,width=10)
        d.ellipse([cx-r*0.72,cy-r*0.72,cx+r*0.72,cy+r*0.72],outline=col,width=9)
        d.ellipse([cx-r*0.28,cy-r*0.28,cx+r*0.28,cy+r*0.28],outline=col,width=9)
    gear(150,980,150); gear(1760,140,120)
    p=ASSET/"bg_grid.png"; im.save(p); return str(p)


# ---------- pptx helpers ----------
def rect(s,x,y,w,h,fill=None,line=None,lw=1.0,shape=MSO_SHAPE.RECTANGLE,radius=None):
    sp=s.shapes.add_shape(shape,Inches(x),Inches(y),Inches(w),Inches(h))
    if radius is not None:
        try: sp.adjustments[0]=radius
        except Exception: pass
    if fill is None: sp.fill.background()
    else: sp.fill.solid(); sp.fill.fore_color.rgb=fill
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb=line; sp.line.width=Pt(lw)
    sp.shadow.inherit=False; return sp

def text(s,x,y,w,h,runs,size=18,color=INK,bold=False,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,
         font=BODY,line=1.05,spacing=2,letter=None):
    tb=s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); tf=tb.text_frame
    tf.word_wrap=True; tf.vertical_anchor=anchor
    for i,para in enumerate(runs if isinstance(runs,list) else [runs]):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment=align; p.line_spacing=line; p.space_after=Pt(spacing); p.space_before=Pt(0)
        for seg in (para if isinstance(para,list) else [(para,color,bold,size,font)]):
            t=seg[0]; r=p.add_run(); r.text=t
            r.font.color.rgb=seg[1] if len(seg)>1 else color
            r.font.bold=seg[2] if len(seg)>2 else bold
            r.font.size=Pt(seg[3] if len(seg)>3 else size)
            r.font.name=seg[4] if len(seg)>4 else font
    return tb

def node(s,x,y,d=0.22):
    rect(s,x,y,d,d,fill=WHITE,line=LINE,lw=1.25,shape=MSO_SHAPE.OVAL)
    rect(s,x+d*0.3,y+d*0.3,d*0.4,d*0.4,fill=CYAN,shape=MSO_SHAPE.OVAL)

def hline(s,x,y,w): rect(s,x,y,w,0.012,fill=LINE)
def vline(s,x,y,h): rect(s,x,y,0.012,h,fill=LINE)

def chevrons(s,x,y,n=3,size=0.34,gap=0.1,color=CYAN):
    for i in range(n):
        rect(s,x+i*(size*0.55+gap),y,size,size,fill=color,shape=MSO_SHAPE.CHEVRON)

def numcircle(s,x,y,n,d=0.62):
    rect(s,x,y,d,d,fill=GREEN,shape=MSO_SHAPE.OVAL)
    text(s,x,y,d,d,[[(str(n),WHITE,True,21,HEAD)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)

def pic_cover(s,path,x,y,w,h):
    """place image cropped to fill the x,y,w,h box (cover)."""
    iw,ih=Image.open(path).size; box=w/h; im=iw/ih
    pic=s.shapes.add_picture(path,Inches(x),Inches(y),Inches(w),Inches(h))
    if im>box:  # too wide -> crop sides
        crop=(1-box/im)/2; pic.crop_left=crop; pic.crop_right=crop
    else:       # too tall -> crop top/bottom
        crop=(1-im/box)/2; pic.crop_top=crop; pic.crop_bottom=crop
    return pic


def main():
    bg=make_bg()
    prs=Presentation(); prs.slide_width=Inches(SW); prs.slide_height=Inches(SH)
    BLANK=prs.slide_layouts[6]
    def newslide():
        s=prs.slides.add_slide(BLANK); s.shapes.add_picture(bg,0,0,Inches(SW),Inches(SH)); return s

    # ===== TITLE =====
    s=newslide()
    # top-left blueprint sensor motif
    hline(s,0,0.9,2.2); node(s,2.05,0.79)
    # right imagery panel (satellite tile) where EMG had the 3D-printer illustration
    pic_cover(s,TV+"nairobi_composites.png",8.35,1.5,4.1,4.1)
    rect(s,8.35,1.5,4.1,4.1,line=TEAL,lw=1.5)
    chevrons(s,12.0,1.0,4,size=0.42,gap=0.0,color=CYAN)  # far-right cyan chevrons (vertical-ish)
    # title
    text(s,0.85,2.0,7.2,2.4,[[("Predicting Village",TEAL,True,52,HEAD)],
                             [("Wealth from Space",TEAL,True,52,HEAD)]],line=1.0)
    text(s,0.9,4.15,3,0.4,[[("P R E S E N T E D   B Y",TEAL_D,True,13,HEAD)]])
    # presenter / context with chevron bullets
    for i,(t,) in enumerate([("Onur Haniffa  ·  231402005",),
                             ("Advisor: Dr. Seda Nilgün Dumlu",),
                             ("ML / DL Internship · Acıbadem MAAÜ · 2026",)]):
        yy=4.6+i*0.62
        chevrons(s,0.95,yy+0.02,3,size=0.22,gap=-0.02,color=CYAN)
        text(s,1.75,yy-0.08,6,0.45,[[(t,INK,False,16,BODY)]])
    # bottom blueprint nodes
    node(s,3.2,6.7); vline(s,3.3,5.9,0.8)
    rect(s,0,7.06,SW,0.012,fill=GREY)

    # ===== STAT SLIDE ("Does it work?") =====
    s=newslide()
    # corner motifs
    node(s,11.9,0.6); hline(s,9.4,0.71,2.5)
    text(s,0,0.75,SW,1.0,[[("Does It Work?  Yes.",TEAL,True,40,HEAD)]],align=PP_ALIGN.CENTER)
    rect(s,SW/2-1.0,1.85,2.0,0.05,fill=CYAN)
    # two big stats
    for x,big,cap in [(2.0,"r 0.78",["Pearson r — matches the","WILDS PovertyMap benchmark"]),
                      (7.4,"0.55",["worst-group r —","beats their 0.45"])]:
        text(s,x,2.4,3.9,1.0,[[(big,TEAL,True,58,HEAD)]],align=PP_ALIGN.CENTER)
        text(s,x,3.7,3.9,1.0,[[(cap[0],INK,False,16,BODY)],[(cap[1],INK,False,16,BODY)]],
             align=PP_ALIGN.CENTER,line=1.2)
    # key takeaway + bullets
    text(s,1.0,4.9,11,0.5,[[("Honest, on unseen countries",TEAL_D,True,18,HEAD)]])
    bullets=["Leave-country-out CV — tested on countries the model never saw",
             "Pooled r² 0.61 on the multi-round (36k-village) data",
             "A naïve random split would have inflated this to 0.73 (+20%) — we refuse that"]
    text(s,1.0,5.45,11,1.6,[[("•  ",CYAN,True,16),(b,INK,False,16,BODY)] for b in bullets],
         line=1.15,spacing=7)

    out="results/poverty_cnn_deck_v2.pptx"; prs.save(out)
    print("saved",out,"with",len(prs.slides._sldIdLst),"slides")


if __name__=="__main__":
    main()
